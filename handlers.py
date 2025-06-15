
import os
import webbrowser
import docx
import openpyxl
import pandas as pd
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image, ImageTk
import tkinter as tk

def display_file_content(app, path):
    try:
        ext = os.path.splitext(path)[1].lower()

        if ext in ['.txt', '.html', '.xml', '.csv']:
            with open(path, 'r', encoding='utf-8') as file:
                content = file.read()
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, content)

        elif ext in ['.doc', '.docx']:
            doc = docx.Document(path)
            content = "".join([para.text for para in doc.paragraphs])
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, content)

        elif ext in ['.xls', '.xlsx']:
            df = pd.read_excel(path, engine='openpyxl')
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, df.to_string())

        elif ext == '.pptx':
            prs = Presentation(path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, content)

        elif ext == '.pdf':
            doc = fitz.open(path)
            content = ""
            for page in doc:
                content += page.get_text()
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, content)

        elif ext in ['.jpeg', '.jpg', '.png', '.gif', '.bmp']:
            image = Image.open(path)
            image.thumbnail((400, 400))
            photo = ImageTk.PhotoImage(image)
            app.image_canvas.delete("all")
            app.image_canvas.create_image(app.image_canvas.winfo_width()//2, app.image_canvas.winfo_height()//2, anchor='center', image=photo)
            app.image_canvas.image = photo

        else:
            app.quiz_text.delete("1.0", tk.END)
            app.quiz_text.insert(tk.END, "Aperçu non disponible pour ce type de fichier.")

    except Exception as e:
        app.quiz_text.delete("1.0", tk.END)
        app.quiz_text.insert(tk.END, f"Erreur lors de l'affichage du fichier : {str(e)}")
