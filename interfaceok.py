import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import webbrowser
from PIL import Image, ImageTk
import docx
import openpyxl
import pandas as pd
from pptx import Presentation
import fitz  # PyMuPDF
import fnmatch
import os
import string
import shutil

from quiz_manager_mod import QuizApp as QuizManagerApp
from file_manager import rename_and_move_files, unzip_file
from quiz_editor import save_quiz, clear_quiz, select_json_file, open_json_editor
from utils import search_files, search_text, open_html_in_browser

class QuizApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Quiz App")
        self.geometry("1200x800")

        style = ttk.Style()
        style.theme_use('clam')

        self.left_frame = tk.Frame(self, width=300, bg='lightgrey')
        self.left_frame.pack(side='left', fill='y')

        self.disk_var = tk.StringVar()
        self.disk_menu = ttk.Combobox(self.left_frame, textvariable=self.disk_var, state='readonly')
        self.disk_menu.pack(pady=10)
        self.disk_menu.bind('<<ComboboxSelected>>', self.on_disk_selected)

        self.search_var = tk.StringVar()
        self.search_entry = ttk.Entry(self.left_frame, textvariable=self.search_var)
        self.search_entry.pack(pady=10)
        self.search_entry.bind('<KeyRelease>', self.on_search)

        self.file_tree_scrollbar = tk.Scrollbar(self.left_frame)
        self.file_tree_scrollbar.pack(side='right', fill='y')

        self.file_tree = ttk.Treeview(self.left_frame, yscrollcommand=self.file_tree_scrollbar.set)
        self.file_tree.pack(pady=10, fill='both', expand=True)
        self.file_tree.bind('<<TreeviewOpen>>', self.on_treeview_open)
        self.file_tree.bind('<<TreeviewSelect>>', self.on_treeview_select)

        self.file_tree_scrollbar.config(command=self.file_tree.yview)

        self.populate_disk_menu()

        self.center_frame = tk.Frame(self, width=600, bg='white')
        self.center_frame.pack(side='left', fill='both', expand=True)

        self.notebook = ttk.Notebook(self.center_frame)
        self.notebook.pack(fill='both', expand=True)

        self.quiz_frame = tk.Frame(self.notebook, bg='white')
        self.notebook.add(self.quiz_frame, text='Éditeur de quiz')

        self.text_scrollbar = tk.Scrollbar(self.quiz_frame)
        self.text_scrollbar.pack(side='right', fill='y')

        self.quiz_text = tk.Text(self.quiz_frame, yscrollcommand=self.text_scrollbar.set)
        self.quiz_text.pack(pady=10, padx=(10, 0), fill='both', expand=True)

        self.text_scrollbar.config(command=self.quiz_text.yview)

        self.save_quiz_button = tk.Button(self.quiz_frame, text="Enregistrer le quiz", command=lambda: save_quiz(self))
        self.save_quiz_button.pack(pady=10)

        self.clear_quiz_button = tk.Button(self.quiz_frame, text="Vider la zone de quiz", command=lambda: clear_quiz(self))
        self.clear_quiz_button.pack(pady=10)

        self.image_frame = tk.Frame(self.notebook, bg='white')
        self.notebook.add(self.image_frame, text='Images')

        self.image_canvas = tk.Canvas(self.image_frame, bg='white')
        self.image_canvas.pack(pady=10, padx=(10, 0), fill='both', expand=True)

        self.html_frame = tk.Frame(self.notebook, bg='white')
        self.notebook.add(self.html_frame, text='HTML')

        self.html_text = tk.Text(self.html_frame)
        self.html_text.pack(pady=10, padx=(10, 0), fill='both', expand=True)

        self.open_html_button = tk.Button(self.html_frame, text="Ouvrir dans le navigateur", command=open_html_in_browser)
        self.open_html_button.pack(pady=10)

        self.open_authorized_html_button = tk.Button(self.html_frame, text="Ouvrir user_management.html", command=lambda: self.open_specific_html("user_management.html"))
        self.open_authorized_html_button.pack(pady=10)

        self.open_authorized_html_button = tk.Button(self.html_frame, text="Ouvrir quiz_user_management.html", command=lambda: self.open_specific_html("quiz_user_management.html"))
        self.open_authorized_html_button.pack(pady=10)

        self.open_authorized_html_button = tk.Button(self.html_frame, text="./fichierquizz", command=lambda: self.open_specific_html("./fichierquizz"))
        self.open_authorized_html_button.pack(pady=10)

        self.open_authorized_html_button = tk.Button(self.html_frame, text="./fichierquizz/manager.html", command=lambda: self.open_specific_html("./fichierquizz/manager.html"))
        self.open_authorized_html_button.pack(pady=10)

        self.open_authorized_html_button = tk.Button(self.html_frame, text="declencheur_auto.html", command=lambda: self.open_specific_html("declencheur_auto.html"))
        self.open_authorized_html_button.pack(pady=10)

        self.lancer_quiz_app_button = tk.Button(self.html_frame, text="Lancer l'application Quiz", command=self.lancer_quiz_app)
        self.lancer_quiz_app_button.pack(pady=10)

        self.load_user_management_html()

        self.right_frame = tk.Frame(self, width=300, bg='lightgrey')
        self.right_frame.pack(side='right', fill='y')

        tk.Button(self.right_frame, text="Renommer et déplacer", command=lambda: rename_and_move_files(self)).pack(pady=5, padx=10, fill='x')
        tk.Button(self.right_frame, text="Décompresser un fichier", command=lambda: unzip_file(self)).pack(pady=5, padx=10, fill='x')
        tk.Button(self.right_frame, text="Sélectionner dossier d'images", command=self.select_image_directory).pack(pady=5, padx=10, fill='x')
        tk.Button(self.right_frame, text="Sélectionner fichier JSON", command=lambda: select_json_file(self)).pack(pady=5, padx=10, fill='x')
        tk.Button(self.right_frame, text="Éditeur JSON", command=lambda: open_json_editor(self)).pack(pady=5, padx=10, fill='x')

        self.progress = ttk.Progressbar(self.right_frame, orient='horizontal', mode='determinate')
        self.progress.pack(pady=10, padx=10, fill='x')

        self.search_file_var = tk.StringVar()
        self.search_file_entry = ttk.Entry(self.right_frame, textvariable=self.search_file_var)
        self.search_file_entry.pack(pady=5, padx=10, fill='x')

        self.search_file_button = tk.Button(self.right_frame, text="Rechercher fichier", command=lambda: search_files(self))
        self.search_file_button.pack(pady=5, padx=10, fill='x')

        self.search_text_var = tk.StringVar()
        self.search_text_entry = ttk.Entry(self.right_frame, textvariable=self.search_text_var)
        self.search_text_entry.pack(pady=5, padx=10, fill='x')

        self.search_text_button = tk.Button(self.right_frame, text="Rechercher dans le texte", command=lambda: search_text(self))
        self.search_text_button.pack(pady=5, padx=10, fill='x')

        self.selected_files = []

        self.quiz_manager_frame = tk.Frame(self.notebook, bg='white')
        self.notebook.add(self.quiz_manager_frame, text='Quiz Manager')
        self.load_quiz_manager_interface()

    def populate_disk_menu(self):
        drives = [f"{d}:\\" for d in string.ascii_uppercase if os.path.exists(f"{d}:\\")]
        self.disk_menu['values'] = drives
        if drives:
            self.disk_menu.current(0)
            self.load_directory(drives[0])

    def on_disk_selected(self, event):
        selected_disk = self.disk_var.get()
        self.load_directory(selected_disk)

    def load_directory(self, path):
        self.file_tree.delete(*self.file_tree.get_children())
        self.insert_node('', path)

    def insert_node(self, parent, path):
        try:
            for item in os.listdir(path):
                item_path = os.path.join(path, item)
                node = self.file_tree.insert(parent, 'end', text=item, values=[item_path])
                if os.path.isdir(item_path):
                    self.file_tree.insert(node, 'end', text='dummy')
        except PermissionError:
            pass

    def on_treeview_open(self, event):
        node = self.file_tree.focus()
        children = self.file_tree.get_children(node)
        if children and self.file_tree.item(children[0], 'text') == 'dummy':
            self.file_tree.delete(children[0])
            path = self.file_tree.item(node, 'values')[0]
            self.insert_node(node, path)

    def on_treeview_select(self, event):
        node = self.file_tree.focus()
        path = self.file_tree.item(node, 'values')[0]
        if os.path.isfile(path):
            self.display_file_content(path)

    def display_file_content(self, path):
        try:
            ext = os.path.splitext(path)[1].lower()
            if ext in ['.txt', '.html', '.xml', '.csv']:
                with open(path, 'r', encoding='utf-8') as file:
                    content = file.read()
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, content)
            elif ext in ['.doc', '.docx']:
                doc = docx.Document(path)
                content = "".join([para.text for para in doc.paragraphs])
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, content)
            elif ext in ['.xls', '.xlsx']:
                df = pd.read_excel(path, engine='openpyxl')
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, df.to_string())
            elif ext == '.pptx':
                prs = Presentation(path)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, content)
            elif ext == '.pdf':
                doc = fitz.open(path)
                content = ""
                for page in doc:
                    content += page.get_text()
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, content)
            elif ext in ['.jpeg', '.jpg', '.png', '.gif', '.bmp']:
                image = Image.open(path)
                image.thumbnail((400, 400))
                photo = ImageTk.PhotoImage(image)
                self.image_canvas.delete("all")
                self.image_canvas.create_image(self.image_canvas.winfo_width()//2, self.image_canvas.winfo_height()//2, anchor='center', image=photo)
                self.image_canvas.image = photo
            else:
                self.quiz_text.delete("1.0", tk.END)
                self.quiz_text.insert(tk.END, "Aperçu non disponible pour ce type de fichier.")
        except Exception as e:
            self.quiz_text.delete("1.0", tk.END)
            self.quiz_text.insert(tk.END, f"Erreur lors de l'affichage du fichier : {str(e)}")

    def on_search(self, event):
        search_term = self.search_var.get().lower()
        for item in self.file_tree.get_children():
            self.file_tree.delete(item)
        self.insert_node('', self.disk_var.get())
        if search_term:
            for item in self.file_tree.get_children():
                if search_term not in self.file_tree.item(item, 'text').lower():
                    self.file_tree.detach(item)

    def select_image_directory(self):
        image_dir = filedialog.askdirectory(title="Sélectionner un répertoire d'images")
        if image_dir:
            messagebox.showinfo("Répertoire sélectionné", f"Répertoire d'images sélectionné : {image_dir}")

    def load_user_management_html(self):
        try:
            with open("user_management.html", "r", encoding="utf-8") as file:
                html_content = file.read()
                self.html_text.delete("1.0", tk.END)
                self.html_text.insert(tk.END, html_content)
        except Exception as e:
            messagebox.showerror("Erreur", f"Erreur lors du chargement du fichier HTML : {str(e)}")

    def open_specific_html(self, filename):
        authorized_path = os.path.abspath(filename)
        if os.path.exists(authorized_path):
            webbrowser.open(f"file://{authorized_path}")
        else:
            print(f"Fichier {filename} non trouvé.")

    def lancer_quiz_app(self):
        import subprocess
        subprocess.Popen(["python", "quiz_manager_modifié.py"])

    def load_quiz_manager_interface(self):
        self.quiz_manager_app = QuizManagerApp(self.quiz_manager_frame)

class MonInterface:
    def __init__(self, root):
        self.root = root
        self.root.title("Mon Application")

        self.notebook = ttk.Notebook(root)
        self.notebook.pack(expand=1, fill="both")

        self.create_other_tabs()

        self.quiz_tab = ttk.Frame(self.notebook)
        self.notebook.add(self.quiz_tab, text="Quiz")
        self.lancer_quiz_app_button = tk.Button(self.quiz_tab, text="Lancer l'application Quiz", command=self.lancer_quiz_app)
        self.lancer_quiz_app_button.pack(pady=10)

        self.lancer_quiz_button = tk.Button(self.quiz_tab, text="Charger les quiz", command=self.lancer_quiz)
        self.lancer_quiz_button.pack(pady=10)

    def create_other_tabs(self):
        self.admin_tab = ttk.Frame(self.notebook)
        self.notebook.add(self.admin_tab, text="Admin")

    def lancer_quiz(self):
        self.quiz_app = QuizManagerApp(self.quiz_tab, "./mes fichiers quiz html")

    def lancer_quiz_app(self):
        import subprocess
        subprocess.Popen(["python", "quiz_manager_modifié.py"])

if __name__ == "__main__":
    
    app = QuizApp()
    app.mainloop()

